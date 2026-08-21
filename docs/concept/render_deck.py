#!/usr/bin/env python3
"""Render a .pptx to PNGs with Pillow, for visual QA where LibreOffice cannot run.

Not a substitute for PowerPoint's own renderer — it approximates text metrics and
ignores effects like shadows. It is enough to see overlap, overflow, alignment
and contrast, which is what the visual QA pass is looking for.

Usage: python3 render_deck.py deck.pptx [outdir] [--dpi 110]
"""
import sys
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

EMU_IN = 914400
FONTS = {
    ("Calibri", 0, 0): "/usr/share/fonts/truetype/crosextra/Carlito-Regular.ttf",
    ("Calibri", 1, 0): "/usr/share/fonts/truetype/crosextra/Carlito-Bold.ttf",
    ("Calibri", 0, 1): "/usr/share/fonts/truetype/crosextra/Carlito-Italic.ttf",
    ("Calibri", 1, 1): "/usr/share/fonts/truetype/crosextra/Carlito-BoldItalic.ttf",
    ("Cambria", 0, 0): "/usr/share/fonts/truetype/crosextra/Caladea-Regular.ttf",
    ("Cambria", 1, 0): "/usr/share/fonts/truetype/crosextra/Caladea-Bold.ttf",
    ("Cambria", 0, 1): "/usr/share/fonts/truetype/crosextra/Caladea-Italic.ttf",
    ("Cambria", 1, 1): "/usr/share/fonts/truetype/crosextra/Caladea-BoldItalic.ttf",
    ("Courier New", 0, 0): "/usr/share/fonts/truetype/liberation/LiberationMono-Regular.ttf",
    ("Courier New", 1, 0): "/usr/share/fonts/truetype/liberation/LiberationMono-Bold.ttf",
}
FALLBACK = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
FALLBACK_B = "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"
_cache = {}


def font(name, pt, bold=False, italic=False, dpi=110):
    px = max(6, int(round(pt * dpi / 72)))
    key = (name, bold, italic, px)
    if key in _cache:
        return _cache[key]
    path = FONTS.get((name, int(bold), int(italic)))
    if path is None or not Path(path).exists():
        path = FONTS.get((name, int(bold), 0))
    if path is None or not Path(path).exists():
        path = FALLBACK_B if bold else FALLBACK
    f = ImageFont.truetype(path, px)
    _cache[key] = f
    return f


def rgb(color_fmt, default=None):
    try:
        if color_fmt and color_fmt.type is not None:
            return "#%02X%02X%02X" % tuple(color_fmt.rgb)
    except Exception:
        pass
    return default


def shape_fill(sh):
    try:
        if sh.fill.type is not None and sh.fill.type == 1:  # solid
            return rgb(sh.fill.fore_color)
    except Exception:
        pass
    return None


def shape_line(sh):
    try:
        c = rgb(sh.line.color)
        w = sh.line.width.pt if sh.line.width else 1.0
        return c, w
    except Exception:
        return None, 1.0


def wrap(draw, text, fnt, max_px):
    out = []
    for hard in text.split("\n"):
        if not hard.strip():
            out.append("")
            continue
        words, cur = hard.split(), ""
        for w in words:
            trial = (cur + " " + w).strip()
            if draw.textlength(trial, font=fnt) <= max_px or not cur:
                cur = trial
            else:
                out.append(cur)
                cur = w
        out.append(cur)
    return out


def draw_text_frame(draw, sh, x, y, w, h, dpi):
    tf = sh.text_frame
    li = (tf.margin_left or 0) / EMU_IN * dpi
    ri = (tf.margin_right or 0) / EMU_IN * dpi
    ti = (tf.margin_top or 0) / EMU_IN * dpi
    bi = (tf.margin_bottom or 0) / EMU_IN * dpi
    box_w = max(4, w - li - ri)

    lines = []
    for para in tf.paragraphs:
        txt = "".join(r.text for r in para.runs)
        pt, name, bold, ital, col = 18.0, "Calibri", False, False, "#16201C"
        for r in para.runs:
            if r.font.size:
                pt = r.font.size.pt
            if r.font.name:
                name = r.font.name
            bold = bold or bool(r.font.bold)
            ital = ital or bool(r.font.italic)
            c = rgb(r.font.color)
            if c:
                col = c
            break
        fnt = font(name, pt, bold, ital, dpi)
        align = para.alignment
        lh = pt * 1.22 * dpi / 72
        if not txt.strip():
            lines.append((None, "", fnt, col, align, lh))
            continue
        for ln in wrap(draw, txt, fnt, box_w):
            lines.append((para, ln, fnt, col, align, lh))

    total = sum(l[5] for l in lines)
    anchor = tf.vertical_anchor
    if anchor == MSO_ANCHOR.MIDDLE:
        cy = y + ti + (h - ti - bi - total) / 2
    elif anchor == MSO_ANCHOR.BOTTOM:
        cy = y + h - bi - total
    else:
        cy = y + ti

    for _, ln, fnt, col, align, lh in lines:
        if ln:
            tw = draw.textlength(ln, font=fnt)
            if align == PP_ALIGN.CENTER:
                tx = x + li + (box_w - tw) / 2
            elif align == PP_ALIGN.RIGHT:
                tx = x + li + box_w - tw
            else:
                tx = x + li
            draw.text((tx, cy), ln, font=fnt, fill=col)
        cy += lh


def render(path, outdir, dpi=110):
    prs = Presentation(path)
    SW = prs.slide_width / EMU_IN
    SH = prs.slide_height / EMU_IN
    Wp, Hp = int(SW * dpi), int(SH * dpi)
    outdir = Path(outdir)
    outdir.mkdir(parents=True, exist_ok=True)
    made = []

    for idx, slide in enumerate(prs.slides, 1):
        bg = "#FFFFFF"
        try:
            if slide.background.fill.type == 1:
                bg = rgb(slide.background.fill.fore_color, "#FFFFFF")
        except Exception:
            pass
        img = Image.new("RGB", (Wp, Hp), bg)
        d = ImageDraw.Draw(img)

        for sh in slide.shapes:
            x = sh.left / EMU_IN * dpi
            y = sh.top / EMU_IN * dpi
            w = (sh.width or 0) / EMU_IN * dpi
            h = (sh.height or 0) / EMU_IN * dpi

            if sh.has_table:
                tbl = sh.table
                colw = [c.width / EMU_IN * dpi for c in tbl.columns]
                rowh = [r.height / EMU_IN * dpi for r in tbl.rows]
                cy = y
                for ri, row in enumerate(tbl.rows):
                    cx = x
                    for ci, cell in enumerate(row.cells):
                        cf = None
                        try:
                            if cell.fill.type == 1:
                                cf = rgb(cell.fill.fore_color)
                        except Exception:
                            pass
                        d.rectangle([cx, cy, cx + colw[ci], cy + rowh[ri]],
                                    fill=cf or "#FFFFFF", outline="#D4DAD2")
                        para = cell.text_frame.paragraphs[0]
                        pt, name, bold, col = 11.0, "Calibri", False, "#16201C"
                        for r in para.runs:
                            if r.font.size:
                                pt = r.font.size.pt
                            if r.font.name:
                                name = r.font.name
                            bold = bool(r.font.bold)
                            col = rgb(r.font.color, col)
                            break
                        fnt = font(name, pt, bold, False, dpi)
                        t = cell.text_frame.text
                        if t:
                            d.text((cx + 6, cy + (rowh[ri] - pt * 1.22 * dpi / 72) / 2),
                                   t, font=fnt, fill=col)
                        cx += colw[ci]
                    cy += rowh[ri]
                continue

            st = sh.shape_type
            fill = shape_fill(sh)
            lc, lw = shape_line(sh)
            lwpx = max(1, int(lw * dpi / 72))

            if st == MSO_SHAPE_TYPE.AUTO_SHAPE or sh.has_text_frame:
                try:
                    name = sh.auto_shape_type
                except Exception:
                    name = None
                sname = str(name) if name else ""
                if "OVAL" in sname or "ELLIPSE" in sname:
                    if fill or lc:
                        d.ellipse([x, y, x + w, y + h], fill=fill,
                                  outline=lc, width=lwpx)
                elif "LINE" in sname and h < 3:
                    if lc:
                        d.line([x, y, x + w, y + h], fill=lc, width=lwpx)
                elif "ARROW" in sname:
                    if fill:
                        if "LEFT" in sname:
                            d.polygon([(x + w, y), (x, y + h / 2), (x + w, y + h)], fill=fill)
                        elif "UP" in sname:
                            d.polygon([(x, y + h), (x + w / 2, y), (x + w, y + h)], fill=fill)
                        elif "DOWN" in sname:
                            d.polygon([(x, y), (x + w / 2, y + h), (x + w, y)], fill=fill)
                        else:
                            d.polygon([(x, y), (x + w, y + h / 2), (x, y + h)], fill=fill)
                elif fill or lc:
                    r = min(10, int(min(w, h) / 4))
                    if r > 2:
                        d.rounded_rectangle([x, y, x + w, y + h], radius=r,
                                            fill=fill, outline=lc, width=lwpx)
                    else:
                        d.rectangle([x, y, x + w, y + h], fill=fill,
                                    outline=lc, width=lwpx)

            if sh.has_text_frame and sh.text_frame.text.strip():
                draw_text_frame(d, sh, x, y, w, h, dpi)

        out = outdir / f"slide-{idx:02d}.png"
        img.save(out)
        made.append(str(out))
    return made


if __name__ == "__main__":
    args = [a for a in sys.argv[1:] if not a.startswith("--")]
    dpi = 110
    for a in sys.argv[1:]:
        if a.startswith("--dpi"):
            dpi = int(a.split("=")[1]) if "=" in a else 110
    deck = args[0] if args else "/home/user/project1/docs/concept/TEA-concept-overview.pptx"
    outdir = args[1] if len(args) > 1 else "/tmp/claude-0/deckpreview"
    for f in render(deck, outdir, dpi):
        print(f)
