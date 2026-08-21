#!/usr/bin/env python3
"""Geometry QA for a generated deck, for environments with no LibreOffice.

LibreOffice cannot render in this sandbox, so the usual look-at-the-slides pass
is unavailable. This checks the defects that pass would catch, against the real
coordinates and text in the written file:

  - shapes outside the slide, or inside the 0.5" edge margin
  - text frames that overlap each other
  - text estimated to overflow its box

Text fit is an estimate from per-family average character widths. Treat a
reported overflow as "look at this", not as proof.
"""
import math
import sys
from pptx import Presentation
from pptx.util import Emu

EMU_IN = 914400
EDGE_MARGIN = 0.5      # inches
OVERLAP_TOL = 0.02     # inches; ignore hairline touching

# average glyph width as a fraction of point size, mixed-case Latin
CHAR_W = {
    "Courier New": 0.600,   # monospace, exact
    "Cambria": 0.497,
    "Calibri": 0.478,
}
DEFAULT_CHAR_W = 0.50


def inches(v):
    return (v or 0) / EMU_IN


def runs_of(shape):
    out = []
    for para in shape.text_frame.paragraphs:
        text = "".join(r.text for r in para.runs)
        size, face, bold = None, None, False
        for r in para.runs:
            if r.font.size and size is None:
                size = r.font.size.pt
            if r.font.name and face is None:
                face = r.font.name
            bold = bold or bool(r.font.bold)
        out.append((text, size or 18.0, face or "Calibri", bold))
    return out


def estimate_lines(text, pt, face, width_in, bold):
    if not text.strip():
        return 1
    cw = CHAR_W.get(face, DEFAULT_CHAR_W) * pt * (1.03 if bold else 1.0)
    usable_pt = width_in * 72 - 8          # allow for text-frame inset
    if usable_pt <= cw:
        return 99
    per_line = max(1, int(usable_pt / cw))
    lines = 0
    for hard in text.split("\n"):
        words, cur = hard.split(), 0
        if not words:
            lines += 1
            continue
        n = 0
        for w in words:
            add = len(w) + (1 if n else 0)
            if n + add > per_line and n:
                lines += 1
                n = len(w)
            else:
                n += add
        lines += 1
    return lines


def check(path):
    prs = Presentation(path)
    SW, SH = inches(prs.slide_width), inches(prs.slide_height)
    problems = []

    for idx, slide in enumerate(prs.slides, 1):
        boxes = []
        for sh in slide.shapes:
            x, y = inches(sh.left), inches(sh.top)
            w, h = inches(sh.width), inches(sh.height)
            has_text = sh.has_text_frame and sh.text_frame.text.strip()
            name = (sh.text_frame.text.strip()[:34].replace("\n", " ")
                    if has_text else sh.shape_type)

            # off-slide
            if x < -0.01 or y < -0.01 or x + w > SW + 0.01 or y + h > SH + 0.01:
                problems.append(
                    f"slide {idx}: off-slide  [{name}]  x={x:.2f} y={y:.2f} "
                    f"w={w:.2f} h={h:.2f}  (slide {SW:.2f}x{SH:.2f})")

            # edge margin — full-bleed backgrounds are intentional, skip them
            full_bleed = w > SW - 0.2 and h > SH - 0.2
            # running footer / page number sit in the margin band by convention
            footer_band = y > SH - 0.75 and h < 0.4
            if has_text and not full_bleed and not footer_band:
                if x < EDGE_MARGIN - 0.01 or y < EDGE_MARGIN - 0.01 \
                   or x + w > SW - EDGE_MARGIN + 0.01 or y + h > SH - EDGE_MARGIN + 0.01:
                    problems.append(
                        f"slide {idx}: inside {EDGE_MARGIN}\" edge margin  [{name}]  "
                        f"x={x:.2f} y={y:.2f} r={x+w:.2f} b={y+h:.2f}")

            # estimated text overflow
            if has_text:
                total = 0.0
                for text, pt, face, bold in runs_of(sh):
                    n = estimate_lines(text, pt, face, w, bold)
                    total += n * pt * 1.22 / 72
                if total > h + 0.06:
                    problems.append(
                        f"slide {idx}: text may overflow  [{name}]  "
                        f"needs ~{total:.2f}\" in {h:.2f}\"")
                boxes.append((x, y, w, h, name))

        # text-frame overlaps
        for i in range(len(boxes)):
            for j in range(i + 1, len(boxes)):
                ax, ay, aw, ah, an = boxes[i]
                bx, by, bw, bh, bn = boxes[j]
                ox = min(ax + aw, bx + bw) - max(ax, bx)
                oy = min(ay + ah, by + bh) - max(ay, by)
                if ox > OVERLAP_TOL and oy > OVERLAP_TOL:
                    problems.append(
                        f"slide {idx}: text overlap  [{an}] x [{bn}]  "
                        f"{ox:.2f}\" x {oy:.2f}\"")

    return problems, len(prs.slides)


if __name__ == "__main__":
    path = sys.argv[1] if len(sys.argv) > 1 else \
        "/home/user/project1/docs/concept/TEA-concept-overview.pptx"
    problems, n = check(path)
    print(f"{n} slides checked")
    if not problems:
        print("no geometry problems found")
        sys.exit(0)
    print(f"{len(problems)} to look at:")
    for p in problems:
        print(f"  {p}")
    sys.exit(1)
